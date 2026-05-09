import json, os, sys, io

_pip = '/home/site/pip-packages'
if os.path.isdir(_pip) and _pip not in sys.path:
    sys.path.insert(0, _pip)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.chart.data import CategoryChartData
except ImportError as e:
    print(f'Error: python-pptx not installed ({e}). sys.path={sys.path}')
    sys.exit(1)

try:
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    import matplotlib.ticker as ticker
    HAS_MPL = True
except ImportError:
    HAS_MPL = False

# -- Color palette (Microsoft Fluent v2 — refined, executive) --
AZURE_BLUE = RGBColor(0x00, 0x67, 0xC0)
AZURE_DEEP = RGBColor(0x00, 0x3F, 0x87)
AZURE_LIGHT = RGBColor(0xEA, 0xF2, 0xFB)
TITLE_BG = RGBColor(0x0B, 0x1F, 0x3A)
TITLE_BG_2 = RGBColor(0x10, 0x35, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLIDE_BG = RGBColor(0xFB, 0xFB, 0xFD)
ACCENT_TEAL = RGBColor(0x00, 0xB7, 0xC3)
ACCENT_GREEN = RGBColor(0x10, 0x7C, 0x10)
ACCENT_AMBER = RGBColor(0xCA, 0x5E, 0x00)
ACCENT_RED = RGBColor(0xC4, 0x32, 0x4C)
TEXT_DARK = RGBColor(0x14, 0x1E, 0x2C)
TEXT_MED = RGBColor(0x4A, 0x52, 0x5C)
TEXT_LIGHT = RGBColor(0x8A, 0x90, 0x99)
RULE_GREY = RGBColor(0xE4, 0xE7, 0xEC)
KPI_BG = RGBColor(0xFF, 0xFF, 0xFF)
KPI_BORDER = RGBColor(0xDC, 0xE2, 0xEA)

CHART_COLORS = ['#0067C0', '#00B7C3', '#107C10', '#CA5E00', '#8661C5',
                '#003F87', '#50E6FF', '#2D7D46', '#E3008C', '#FFB900',
                '#B4009E', '#0B1F3A', '#4F6BED', '#C239B3', '#767676']

slides_data = json.loads(sys.stdin.read())
output_path = os.environ['PPTX_OUTPUT_PATH']
customer_name = os.environ.get('PPTX_CUSTOMER', '').strip()
total_slides = len(slides_data)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=TEXT_DARK, alignment=PP_ALIGN.LEFT, font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_brand_footer(slide, slide_num, total, customer):
    rule = slide.shapes.add_shape(1, Inches(0.5), Inches(SLIDE_H_IN - 0.45), Inches(SLIDE_W_IN - 1.0), Inches(0.012))
    rule.fill.solid(); rule.fill.fore_color.rgb = RULE_GREY
    rule.line.fill.background()
    mark = slide.shapes.add_shape(1, Inches(0.5), Inches(SLIDE_H_IN - 0.32), Inches(0.14), Inches(0.14))
    mark.fill.solid(); mark.fill.fore_color.rgb = AZURE_BLUE
    mark.line.fill.background()
    left_label = 'Microsoft Azure FinOps'
    if customer:
        left_label = f'{customer}  \u2022  Microsoft Azure FinOps'
    tb = slide.shapes.add_textbox(Inches(0.72), Inches(SLIDE_H_IN - 0.36), Inches(8.5), Inches(0.25))
    p = tb.text_frame.paragraphs[0]
    p.text = left_label
    p.font.size = Pt(9); p.font.name = 'Segoe UI'; p.font.color.rgb = TEXT_MED
    p.font.bold = True
    tb2 = slide.shapes.add_textbox(Inches(SLIDE_W_IN - 1.5), Inches(SLIDE_H_IN - 0.36), Inches(1.0), Inches(0.25))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = f'{slide_num} / {total}'
    p2.font.size = Pt(9); p2.font.name = 'Segoe UI'; p2.font.color.rgb = TEXT_LIGHT
    p2.alignment = PP_ALIGN.RIGHT

def apply_picture_shadow(picture):
    """Add a subtle drop shadow to a picture via XML so chart images feel lifted."""
    try:
        from pptx.oxml.ns import qn
        from lxml import etree
        spPr = picture._element.spPr
        for child in spPr.findall(qn('a:effectLst')):
            spPr.remove(child)
        eff = etree.SubElement(spPr, qn('a:effectLst'))
        outer = etree.SubElement(eff, qn('a:outerShdw'),
                                 {'blurRad': '50800', 'dist': '25400', 'dir': '5400000',
                                  'algn': 'tl', 'rotWithShape': '0'})
        srgb = etree.SubElement(outer, qn('a:srgbClr'), {'val': '141E2C'})
        etree.SubElement(srgb, qn('a:alpha'), {'val': '18000'})
    except Exception:
        pass

def _fit_bullets_pt(n):
    """Pick a bullet font size that scales inversely with bullet count.\n    Few bullets read big and confident; many bullets stay legible."""
    if n <= 3: return 24
    if n <= 5: return 20
    if n <= 7: return 17
    if n <= 10: return 15
    return 13

def _fit_kpi_value_pt(text, card_width_in):
    """Pick a value font size that fits a single line inside the KPI card.
    Heuristic: ~0.55in per character at 32pt; scale down by length and card width."""
    n = len(str(text))
    # Available chars at 32pt baseline given card width (minus padding)
    avail_in = max(0.5, card_width_in - 0.5)
    # Approximate avg char width in inches at given pt size: pt * 0.0095 (Segoe UI bold)
    for size in (32, 28, 24, 20, 18, 16, 14):
        if n * size * 0.0095 <= avail_in:
            return size
    return 14

def add_kpi_card(slide, left, top, width, height, label, value, sublabel='', accent=AZURE_BLUE):
    """A single executive KPI tile: rounded card, thin colored top bar, big value, small label below."""
    card = slide.shapes.add_shape(5, Inches(left), Inches(top), Inches(width), Inches(height))
    card.adjustments[0] = 0.08
    card.fill.solid(); card.fill.fore_color.rgb = KPI_BG
    card.line.color.rgb = KPI_BORDER
    card.line.width = Pt(0.75)
    bar = slide.shapes.add_shape(1, Inches(left + 0.12), Inches(top + 0.12), Inches(width - 0.24), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_textbox(slide, left + 0.25, top + 0.3, width - 0.5, 0.35,
                label.upper(), font_size=10, bold=True, color=TEXT_LIGHT)
    value_pt = _fit_kpi_value_pt(value, width)
    val_box = add_textbox(slide, left + 0.25, top + 0.7, width - 0.5, 0.9,
                value, font_size=value_pt, bold=True, color=TEXT_DARK)
    # Single-line: disable word wrap so the chosen font size is honored without runaway wrap
    val_box.text_frame.word_wrap = False
    if sublabel:
        add_textbox(slide, left + 0.25, top + 1.6, width - 0.5, 0.45,
                    sublabel, font_size=11, color=TEXT_MED)

def stars_text(score):
    s = max(0, min(5, int(round(float(score)))))
    return '★' * s + '☆' * (5 - s)

def star_color(score):
    s = float(score)
    if s >= 4: return ACCENT_GREEN
    if s >= 3: return AZURE_BLUE
    if s >= 2: return ACCENT_AMBER
    return ACCENT_RED

def add_maturity_table(slide, left, top, width, height, rows):
    """rows = [{label, before, after, action}] — renders a clean before/after grid with star ratings."""
    n_rows = len(rows) + 1  # +header
    n_cols = 4
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = tbl_shape.table
    # column widths: dimension | before | after | action
    col_widths = [3.2, 1.6, 1.6, width - 6.4]
    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)
    headers = ['Dimension', 'Before', 'After', 'What we did']
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = ''
        cell.fill.solid(); cell.fill.fore_color.rgb = TITLE_BG
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = 'Segoe UI'
        cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.12)
        cell.margin_top = Inches(0.06); cell.margin_bottom = Inches(0.06)
    for r, row in enumerate(rows, start=1):
        zebra = SLIDE_BG if r % 2 == 0 else WHITE
        for c in range(n_cols):
            cell = table.cell(r, c)
            cell.fill.solid(); cell.fill.fore_color.rgb = zebra
            cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.12)
            cell.margin_top = Inches(0.05); cell.margin_bottom = Inches(0.05)
        # Dimension name
        p0 = table.cell(r, 0).text_frame.paragraphs[0]
        p0.text = row.get('label', ''); p0.font.size = Pt(11); p0.font.bold = True; p0.font.color.rgb = TEXT_DARK; p0.font.name = 'Segoe UI'
        # Before stars
        before = row.get('before', 0)
        p1 = table.cell(r, 1).text_frame.paragraphs[0]
        p1.text = stars_text(before); p1.font.size = Pt(14); p1.font.color.rgb = star_color(before); p1.font.name = 'Segoe UI'
        # After stars
        after = row.get('after', 0)
        p2 = table.cell(r, 2).text_frame.paragraphs[0]
        p2.text = stars_text(after); p2.font.size = Pt(14); p2.font.color.rgb = star_color(after); p2.font.name = 'Segoe UI'
        # Action description
        p3 = table.cell(r, 3).text_frame.paragraphs[0]
        p3.text = row.get('action', ''); p3.font.size = Pt(10); p3.font.color.rgb = TEXT_MED; p3.font.name = 'Segoe UI'

def add_bullets(slide, left, top, width, height, items, font_size=14, color=TEXT_DARK):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Segoe UI'
        p.space_after = Pt(6)
        p.level = 0
    return txBox

def render_chart_image(chart_cfg, path):
    if not HAS_MPL:
        return False
    ctype = chart_cfg.get('type', 'bar')
    labels = chart_cfg.get('labels', [])
    values = chart_cfg.get('values', [])
    if not labels or not values:
        return False
    title = chart_cfg.get('title', '')
    colors = chart_cfg.get('colors', CHART_COLORS[:len(labels)])
    if len(colors) < len(labels):
        colors = colors + CHART_COLORS[len(colors):len(labels)]

    if ctype == 'pie':
        fig, ax = plt.subplots(figsize=(10, 5))
    elif ctype == 'horizontal_bar':
        fig_h = max(4, len(labels) * 0.55)
        fig, ax = plt.subplots(figsize=(8, fig_h))
    elif len(labels) > 15:
        fig, ax = plt.subplots(figsize=(10, 4.5))
    else:
        fig, ax = plt.subplots(figsize=(8, 4.5))
    fig.patch.set_facecolor('#FAFAFA')
    ax.set_facecolor('#FAFAFA')

    if ctype == 'pie':
        ax.set_aspect('equal')
        total = sum(values)
        threshold = 0.03 * total if total > 0 else 0
        main_labels, main_values, main_colors = [], [], []
        other_val = 0
        for lbl, val, clr in zip(labels, values, colors):
            if val >= threshold:
                main_labels.append(lbl); main_values.append(val); main_colors.append(clr)
            else:
                other_val += val
        if other_val > 0:
            main_labels.append('Other'); main_values.append(other_val); main_colors.append('#D0D0D0')

        wedges, autotexts = ax.pie(main_values, colors=main_colors,
                                    autopct='', startangle=90,
                                    pctdistance=0.8, wedgeprops={'width': 0.55, 'edgecolor': 'white', 'linewidth': 2})[0:2]
        centre = plt.Circle((0, 0), 0.45, fc='#FAFAFA')
        ax.add_artist(centre)
        ax.text(0, 0, f'${total:,.0f}', ha='center', va='center', fontsize=16, fontweight='bold', color='#323130')
        ax.text(0, -0.15, 'Total (USD)', ha='center', va='center', fontsize=9, color='#605E5C')
        legend_labels = [f'{lbl}  ${val:,.0f} ({val/total*100:.1f}%)' for lbl, val in zip(main_labels, main_values)]
        ax.legend(wedges, legend_labels, loc='center left', bbox_to_anchor=(1.05, 0.5),
                  fontsize=9, frameon=False, labelspacing=1.0)
    elif ctype == 'horizontal_bar':
        y_pos = range(len(labels))
        bars = ax.barh(y_pos, values, color=colors[:len(labels)], height=0.6, edgecolor='white', linewidth=0.5)
        ax.set_yticks(y_pos)
        ax.set_yticklabels(labels, fontsize=10, color='#323130')
        ax.invert_yaxis()
        for bar, val in zip(bars, values):
            ax.text(bar.get_width() + max(values)*0.02, bar.get_y() + bar.get_height()/2,
                    f'${val:,.0f}', va='center', fontsize=9, color='#605E5C')
    elif ctype == 'line':
        line_color = colors[0] if colors else '#0078D4'
        x_pos = range(len(labels))
        ax.plot(x_pos, values, color=line_color,
                marker='o', linewidth=2.5, markersize=5, markerfacecolor='white', markeredgewidth=2, markeredgecolor=line_color)
        ax.fill_between(x_pos, values, alpha=0.12, color=line_color)
        ax.set_xticks(x_pos)
        ax.set_xticklabels(labels, rotation=45, ha='right', fontsize=9, color='#605E5C')
        ax.set_xlim(-0.5, len(labels) - 0.5)
    elif ctype == 'waterfall':
        # Waterfall: first bar = baseline (neutral), middle bars = deltas (green +, red -), last = total (blue).
        # `values` is the array of deltas; first item = starting baseline, last item = final total (auto-computed if 0).
        baseline = float(values[0]) if values else 0
        deltas = [float(v) for v in values[1:-1]] if len(values) > 2 else []
        running = baseline
        positions = [0]
        bar_vals = [baseline]
        bar_colors = ['#4A525C']
        bottoms = [0]
        for d in deltas:
            positions.append(len(positions))
            bottoms.append(running if d >= 0 else running + d)
            bar_vals.append(abs(d))
            bar_colors.append('#107C10' if d >= 0 else '#C4324C')
            running += d
        positions.append(len(positions))
        final_total = float(values[-1]) if len(values) > 1 and values[-1] != 0 else running
        bar_vals.append(final_total)
        bar_colors.append('#0067C0')
        bottoms.append(0)
        bars = ax.bar(positions, bar_vals, bottom=bottoms, color=bar_colors,
                      width=0.6, edgecolor='white', linewidth=1)
        ax.set_xticks(positions)
        ax.set_xticklabels(labels, fontsize=10, color='#323130', rotation=20, ha='right')
        # Value labels above each bar
        all_deltas_padded = [baseline] + deltas + [final_total]
        for i, (b, v, btm) in enumerate(zip(bars, bar_vals, bottoms)):
            top = btm + v
            label_val = all_deltas_padded[i]
            sign = '+' if 0 < i < len(positions) - 1 and label_val >= 0 else ('' if i == 0 or i == len(positions) - 1 else '−')
            display = f'{sign}{abs(label_val):,.0f}' if i not in (0, len(positions) - 1) else f'{label_val:,.0f}'
            ax.text(b.get_x() + b.get_width() / 2, top + max(bar_vals) * 0.02,
                    display, ha='center', va='bottom', fontsize=10, fontweight='600', color='#323130')
        ax.set_ylim(0, max([b + h for b, h in zip(bottoms, bar_vals)]) * 1.15)
    else:
        bars = ax.bar(labels, values, color=colors[:len(labels)], width=0.65, edgecolor='white', linewidth=0.5)
        plt.xticks(rotation=45, ha='right', fontsize=9, color='#605E5C')
        for bar, val in zip(bars, values):
            ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + max(values)*0.01,
                    f'${val:,.0f}', ha='center', va='bottom', fontsize=9, color='#605E5C')

    if title:
        ax.set_title(title, fontsize=14, fontweight='600', pad=15, color='#323130')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#E8E8E8')
    ax.spines['bottom'].set_color('#E8E8E8')
    ax.tick_params(colors='#605E5C', labelsize=9)
    if ctype == 'horizontal_bar':
        ax.xaxis.set_major_formatter(ticker.FuncFormatter(lambda x, p: f'${x:,.0f}'))
    elif ctype not in ('pie', 'waterfall'):
        ax.yaxis.set_major_formatter(ticker.FuncFormatter(lambda x, p: f'${x:,.0f}'))
    plt.tight_layout()
    fig.savefig(path, dpi=150, bbox_inches='tight', facecolor='#FAFAFA')
    plt.close(fig)
    return True

for idx, slide_def in enumerate(slides_data):
    layout = slide_def.get('layout', 'content')
    title = slide_def.get('title', '')
    subtitle = slide_def.get('subtitle', '')
    bullets = slide_def.get('bullets', [])
    chart_cfg = slide_def.get('chart', None)
    notes = slide_def.get('notes', '')

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    if layout == 'title':
        add_bg(slide, TITLE_BG)
        # Two-tone diagonal accent block — gives the title slide depth
        accent_block = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(SLIDE_W_IN * 0.55), prs.slide_height)
        accent_block.fill.solid(); accent_block.fill.fore_color.rgb = TITLE_BG_2
        accent_block.line.fill.background()
        # Thin teal bar at the very top
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.08))
        shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT_TEAL
        shape.line.fill.background()
        # Brand mark dot
        dot = slide.shapes.add_shape(9, Inches(1), Inches(1.4), Inches(0.32), Inches(0.32))
        dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT_TEAL
        dot.line.fill.background()
        add_textbox(slide, 1.5, 1.35, 11, 0.45, 'AZURE FINOPS', font_size=12, bold=True,
                    color=ACCENT_TEAL)
        # Constrain title within the dark accent block (~6.3in wide) and shrink for long strings
        title_len = len(title)
        if title_len <= 28:
            title_pt = 46
        elif title_len <= 40:
            title_pt = 38
        elif title_len <= 55:
            title_pt = 30
        else:
            title_pt = 24
        add_textbox(slide, 1, 2.4, 6.3, 2.4, title, font_size=title_pt, bold=True, color=WHITE)
        if subtitle:
            add_textbox(slide, 1, 4.7, 6.3, 0.8, subtitle, font_size=18, color=RGBColor(0xBE, 0xD7, 0xEF))
        # Hairline divider
        rule = slide.shapes.add_shape(1, Inches(1), Inches(5.7), Inches(5), Inches(0.012))
        rule.fill.solid(); rule.fill.fore_color.rgb = ACCENT_TEAL
        rule.line.fill.background()
        brand_label = 'Microsoft Azure FinOps Agent'
        if customer_name:
            # Truncate very long customer names so the footer line stays inside the accent block
            display_customer = customer_name if len(customer_name) <= 38 else customer_name[:35] + '…'
            brand_label = f'Prepared for {display_customer}'
        add_textbox(slide, 1, 5.9, 6.3, 0.4, brand_label, font_size=12, bold=True, color=WHITE)
        add_textbox(slide, 1, 6.35, 6.3, 0.4, 'Generated live from Azure APIs', font_size=10, color=RGBColor(0x80, 0x9D, 0xB8))

    elif layout == 'section':
        add_bg(slide, AZURE_LIGHT)
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.18), prs.slide_height)
        shape.fill.solid(); shape.fill.fore_color.rgb = AZURE_BLUE
        shape.line.fill.background()
        eyebrow = slide_def.get('eyebrow', '')
        if eyebrow:
            add_textbox(slide, 1, 2.2, 11, 0.4, eyebrow.upper(), font_size=12, bold=True, color=AZURE_DEEP)
        add_textbox(slide, 1, 2.7, 11, 1.8, title, font_size=44, bold=True, color=TITLE_BG)
        if subtitle:
            add_textbox(slide, 1, 4.5, 11, 0.8, subtitle, font_size=18, color=RGBColor(0x00, 0x5A, 0x9E))

    elif layout == 'chart':
        add_bg(slide, SLIDE_BG)
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.06))
        shape.fill.solid(); shape.fill.fore_color.rgb = AZURE_BLUE
        shape.line.fill.background()
        add_textbox(slide, 0.8, 0.3, 11, 0.7, title, font_size=24, bold=True, color=TEXT_DARK)

        if chart_cfg and HAS_MPL:
            chart_path = output_path.replace('.pptx', f'_chart_{idx}.png')
            if render_chart_image(chart_cfg, chart_path):
                pic = slide.shapes.add_picture(chart_path, Inches(1.5), Inches(1.3), Inches(8), Inches(4.5))
                apply_picture_shadow(pic)
                os.remove(chart_path)
        if bullets:
            bullet_top = 1.3 if not chart_cfg else 6.0
            # Bigger takeaway font when there are only 1–2 lines below a chart
            takeaway_pt = 16 if len(bullets) <= 2 else 14
            add_bullets(slide, 0.8, bullet_top, 11.7, 1.2, bullets, font_size=takeaway_pt, color=TEXT_MED)

    elif layout == 'kpi':
        add_bg(slide, SLIDE_BG)
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.06))
        shape.fill.solid(); shape.fill.fore_color.rgb = AZURE_BLUE
        shape.line.fill.background()
        add_textbox(slide, 0.8, 0.3, 11, 0.7, title, font_size=24, bold=True, color=TEXT_DARK)
        if subtitle:
            add_textbox(slide, 0.8, 0.95, 11, 0.4, subtitle, font_size=13, color=TEXT_MED)
        kpis = slide_def.get('kpis', [])[:4]
        if kpis:
            n = len(kpis)
            gutter = 0.3
            margin = 0.8
            avail = SLIDE_W_IN - 2 * margin - (n - 1) * gutter
            card_w = avail / n
            card_h = 2.1
            top = 1.9
            accent_cycle = [AZURE_BLUE, ACCENT_TEAL, ACCENT_GREEN, ACCENT_AMBER]
            for i, k in enumerate(kpis):
                left = margin + i * (card_w + gutter)
                accent_name = (k.get('accent') or '').lower()
                accent_map = {'blue': AZURE_BLUE, 'teal': ACCENT_TEAL, 'green': ACCENT_GREEN,
                              'amber': ACCENT_AMBER, 'red': ACCENT_RED}
                accent = accent_map.get(accent_name, accent_cycle[i % len(accent_cycle)])
                add_kpi_card(slide, left, top, card_w, card_h,
                             k.get('label', ''), k.get('value', ''), k.get('sublabel', ''), accent)
        if bullets:
            add_bullets(slide, 0.8, 4.3, 11.7, 2.5, bullets, font_size=14, color=TEXT_DARK)

    elif layout == 'maturity':
        add_bg(slide, SLIDE_BG)
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.06))
        shape.fill.solid(); shape.fill.fore_color.rgb = AZURE_BLUE
        shape.line.fill.background()
        add_textbox(slide, 0.8, 0.3, 11, 0.7, title, font_size=24, bold=True, color=TEXT_DARK)
        if subtitle:
            add_textbox(slide, 0.8, 0.95, 11.7, 0.4, subtitle, font_size=13, color=TEXT_MED)
        rows = slide_def.get('rows', [])
        if rows:
            tbl_top = 1.5 if subtitle else 1.2
            tbl_h = max(2.5, len(rows) * 0.45 + 0.5)
            add_maturity_table(slide, 0.8, tbl_top, SLIDE_W_IN - 1.6, tbl_h, rows)
        # Optional summary line under the table
        footnote = slide_def.get('footnote', '')
        if footnote:
            add_textbox(slide, 0.8, SLIDE_H_IN - 0.85, SLIDE_W_IN - 1.6, 0.35,
                        footnote, font_size=12, bold=True, color=AZURE_DEEP)

    elif layout == 'two_column':
        add_bg(slide, SLIDE_BG)
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.06))
        shape.fill.solid(); shape.fill.fore_color.rgb = AZURE_BLUE
        shape.line.fill.background()
        add_textbox(slide, 0.8, 0.3, 11, 0.7, title, font_size=24, bold=True, color=TEXT_DARK)
        right_bullets = slide_def.get('bullets_right', [])
        max_n = max(len(bullets or []), len(right_bullets or []))
        col_pt = _fit_bullets_pt(max_n)
        if bullets:
            add_bullets(slide, 0.8, 1.3, 5.6, 5.5, bullets, font_size=col_pt, color=TEXT_DARK)
        if right_bullets:
            add_bullets(slide, 6.9, 1.3, 5.6, 5.5, right_bullets, font_size=col_pt, color=TEXT_DARK)

    else:
        add_bg(slide, SLIDE_BG)
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.06))
        shape.fill.solid(); shape.fill.fore_color.rgb = AZURE_BLUE
        shape.line.fill.background()
        add_textbox(slide, 0.8, 0.3, 11, 0.7, title, font_size=24, bold=True, color=TEXT_DARK)
        if subtitle:
            add_textbox(slide, 0.8, 1.0, 11, 0.5, subtitle, font_size=16, color=TEXT_MED)
        if bullets:
            bullet_top = 1.7 if subtitle else 1.4
            # Scale bullet font with bullet count so few bullets read big and many still fit
            bullet_pt = _fit_bullets_pt(len(bullets))
            add_bullets(slide, 0.8, bullet_top, 11.7, 5.2, bullets, font_size=bullet_pt, color=TEXT_DARK)

        if chart_cfg and HAS_MPL:
            chart_path = output_path.replace('.pptx', f'_chart_{idx}.png')
            if render_chart_image(chart_cfg, chart_path):
                chart_top = 1.3 + len(bullets) * 0.35 if bullets else 1.3
                pic = slide.shapes.add_picture(chart_path, Inches(1.5), Inches(min(chart_top, 3.0)), Inches(8), Inches(4.0))
                apply_picture_shadow(pic)
                os.remove(chart_path)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    if layout != 'title':
        add_brand_footer(slide, idx + 1, total_slides, customer_name)

prs.save(output_path)
print(f'OK:{output_path}')
print(f'SLIDES:{len(slides_data)}')
